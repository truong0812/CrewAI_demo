"""
Công cụ tạo file PowerPoint từ structured content.
Hỗ trợ nhiều theme, formatting tự động.
"""

import json
import os
from pathlib import Path

from pptx import Presentation
from pptx.util import Inches, Pt, Emu
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN, MSO_ANCHOR
from pptx.enum.shapes import MSO_SHAPE

from config import get_theme


class PPTXGenerator:
    """Tạo file PowerPoint từ JSON outline."""

    def __init__(self, output_dir: str = "output"):
        self.output_dir = output_dir
        os.makedirs(output_dir, exist_ok=True)

    def _hex_to_rgb(self, hex_color: str) -> RGBColor:
        """Chuyển đổi hex color sang RGBColor."""
        hex_color = hex_color.lstrip("#")
        return RGBColor(
            int(hex_color[0:2], 16),
            int(hex_color[2:4], 16),
            int(hex_color[4:6], 16),
        )

    def _add_background(self, slide, bg_color: str):
        """Thêm background color cho slide."""
        background = slide.background
        fill = background.fill
        fill.solid()
        fill.fore_color.rgb = self._hex_to_rgb(bg_color)

    def _add_accent_bar(self, slide, accent_color: str, left=0, top=0, width=Inches(0.15), height=None):
        """Thêm thanh accent color ở bên trái slide."""
        if height is None:
            height = Inches(7.5)
        shape = slide.shapes.add_shape(
            MSO_SHAPE.RECTANGLE, left, top, width, height
        )
        shape.fill.solid()
        shape.fill.fore_color.rgb = self._hex_to_rgb(accent_color)
        shape.line.fill.background()
        return shape

    def _add_title_slide(self, prs, slide_data: dict, theme: dict):
        """Tạo slide tiêu đề."""
        slide = prs.slides.add_slide(prs.slide_layouts[6])  # Blank layout
        self._add_background(slide, theme["bg_color"])
        self._add_accent_bar(slide, theme["accent_color"])

        # Title
        title_box = slide.shapes.add_textbox(
            Inches(1.0), Inches(2.0), Inches(8.0), Inches(2.0)
        )
        tf = title_box.text_frame
        tf.word_wrap = True
        p = tf.paragraphs[0]
        p.text = slide_data.get("title", "")
        p.font.size = Pt(40)
        p.font.bold = True
        p.font.color.rgb = self._hex_to_rgb(theme["title_color"])
        p.alignment = PP_ALIGN.LEFT

        # Subtitle
        if slide_data.get("subtitle"):
            sub_box = slide.shapes.add_textbox(
                Inches(1.0), Inches(4.2), Inches(8.0), Inches(1.0)
            )
            tf = sub_box.text_frame
            tf.word_wrap = True
            p = tf.paragraphs[0]
            p.text = slide_data["subtitle"]
            p.font.size = Pt(20)
            p.font.color.rgb = self._hex_to_rgb(theme["text_color"])
            p.alignment = PP_ALIGN.LEFT

        # Speaker notes
        if slide_data.get("notes"):
            notes_slide = slide.notes_slide
            notes_slide.notes_text_frame.text = slide_data["notes"]

        return slide

    def _add_content_slide(self, prs, slide_data: dict, theme: dict, slide_index: int):
        """Tạo slide nội dung."""
        slide = prs.slides.add_slide(prs.slide_layouts[6])  # Blank layout
        self._add_background(slide, theme["bg_color"])
        self._add_accent_bar(slide, theme["accent_color"])

        # Title
        title_box = slide.shapes.add_textbox(
            Inches(0.8), Inches(0.4), Inches(8.5), Inches(1.0)
        )
        tf = title_box.text_frame
        tf.word_wrap = True
        p = tf.paragraphs[0]
        p.text = slide_data.get("title", f"Slide {slide_index}")
        p.font.size = Pt(28)
        p.font.bold = True
        p.font.color.rgb = self._hex_to_rgb(theme["title_color"])
        p.alignment = PP_ALIGN.LEFT

        # Divider line
        divider = slide.shapes.add_shape(
            MSO_SHAPE.RECTANGLE,
            Inches(0.8), Inches(1.4), Inches(8.5), Pt(2),
        )
        divider.fill.solid()
        divider.fill.fore_color.rgb = self._hex_to_rgb(theme["accent_color"])
        divider.line.fill.background()

        # Content (bullet points)
        content_box = slide.shapes.add_textbox(
            Inches(1.0), Inches(1.7), Inches(8.3), Inches(5.3)
        )
        tf = content_box.text_frame
        tf.word_wrap = True

        bullet_points = slide_data.get("bullet_points", [])
        for i, point in enumerate(bullet_points):
            if i == 0:
                p = tf.paragraphs[0]
            else:
                p = tf.add_paragraph()

            # Xử lý indent level
            if isinstance(point, dict):
                text = point.get("text", "")
                level = point.get("level", 0)
            else:
                text = str(point)
                level = 0

            p.text = text
            p.level = level
            p.font.size = Pt(18) if level == 0 else Pt(15)
            p.font.color.rgb = self._hex_to_rgb(theme["text_color"])
            p.space_after = Pt(8)

            # Bullet character
            if level == 0:
                p.text = "• " + p.text
            else:
                p.text = "  ◦ " + p.text

        # Speaker notes
        if slide_data.get("notes"):
            notes_slide = slide.notes_slide
            notes_slide.notes_text_frame.text = slide_data["notes"]

        return slide

    def _add_summary_slide(self, prs, slide_data: dict, theme: dict):
        """Tạo slide tóm tắt/cảm ơn."""
        slide = prs.slides.add_slide(prs.slide_layouts[6])  # Blank layout
        self._add_background(slide, theme["bg_color"])
        self._add_accent_bar(slide, theme["accent_color"])

        # Title
        title_box = slide.shapes.add_textbox(
            Inches(1.0), Inches(2.5), Inches(8.0), Inches(1.5)
        )
        tf = title_box.text_frame
        tf.word_wrap = True
        p = tf.paragraphs[0]
        p.text = slide_data.get("title", "Cảm ơn!")
        p.font.size = Pt(40)
        p.font.bold = True
        p.font.color.rgb = self._hex_to_rgb(theme["title_color"])
        p.alignment = PP_ALIGN.CENTER

        # Subtitle/content
        bullet_points = slide_data.get("bullet_points", [])
        if bullet_points:
            content_box = slide.shapes.add_textbox(
                Inches(1.5), Inches(4.0), Inches(7.0), Inches(2.5)
            )
            tf = content_box.text_frame
            tf.word_wrap = True
            for i, point in enumerate(bullet_points):
                if i == 0:
                    p = tf.paragraphs[0]
                else:
                    p = tf.add_paragraph()
                text = point if isinstance(point, str) else point.get("text", "")
                p.text = "• " + text
                p.font.size = Pt(18)
                p.font.color.rgb = self._hex_to_rgb(theme["text_color"])
                p.space_after = Pt(6)

        return slide

    def generate(
        self,
        outline_json: str,
        theme_name: str = "corporate",
        filename: str = None,
    ) -> str:
        """
        Tạo file PowerPoint từ JSON outline.

        Args:
            outline_json: JSON string chứa cấu trúc slide
            theme_name: Tên theme (corporate, creative, academic, technical, nature)
            filename: Tên file output (mặc định tự tạo)

        Returns:
            Đường dẫn file đã tạo

        Expected JSON format:
        {
            "presentation_title": "Tiêu đề bài thuyết trình",
            "slides": [
                {
                    "type": "title",
                    "title": "Tiêu đề chính",
                    "subtitle": "Phụ đề"
                },
                {
                    "type": "content",
                    "title": "Tiêu đề slide",
                    "bullet_points": ["Điểm 1", "Điểm 2", {"text": "Sub-point", "level": 1}],
                    "notes": "Speaker notes"
                },
                {
                    "type": "summary",
                    "title": "Tóm tắt / Cảm ơn",
                    "bullet_points": ["Tóm tắt 1", "Tóm tắt 2"]
                }
            ]
        }
        """
        # Parse JSON
        if isinstance(outline_json, str):
            outline = json.loads(outline_json)
        else:
            outline = outline_json

        theme = get_theme(theme_name)

        # Tạo presentation
        prs = Presentation()
        prs.slide_width = Inches(10)
        prs.slide_height = Inches(7.5)

        slides = outline.get("slides", [])

        for i, slide_data in enumerate(slides):
            slide_type = slide_data.get("type", "content")

            if slide_type == "title":
                self._add_title_slide(prs, slide_data, theme)
            elif slide_type == "summary":
                self._add_summary_slide(prs, slide_data, theme)
            else:
                self._add_content_slide(prs, slide_data, theme, i + 1)

        # Nếu không có slide nào, tạo slide mặc định
        if not slides:
            self._add_title_slide(
                prs,
                {
                    "title": outline.get("presentation_title", "Thuyết trình"),
                    "subtitle": "Nội dung đang được xây dựng...",
                },
                theme,
            )

        # Tạo filename
        if filename is None:
            safe_title = "".join(
                c for c in outline.get("presentation_title", "presentation") if c.isalnum() or c in " -_"
            ).strip()[:50]
            filename = f"{safe_title}.pptx"

        if not filename.endswith(".pptx"):
            filename += ".pptx"

        filepath = os.path.join(self.output_dir, filename)
        prs.save(filepath)

        return filepath